"""
RAG Core - Qdrant 版本 (LM Studio 后端，Qdrant 向量库)
"""
import os
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

import glob
import hashlib
import json
from datetime import datetime
from typing import List, Dict, Optional, Any
from dataclasses import dataclass
import requests
import pymupdf
import time

# Qdrant 客户端
from qdrant_client import QdrantClient
from qdrant_client.http.models import (
    Distance,
    VectorParams,
    PointStruct,
    Filter,
    FieldCondition,
    MatchValue
)

@dataclass
class QueryResult:
    answer: str
    sources: List[Dict[str, Any]]
    query: str
    timestamp: str
    # Token 统计和耗时
    prompt_tokens: int = 0
    completion_tokens: int = 0
    total_tokens: int = 0
    query_time: float = 0.0

@dataclass
class IndexStatus:
    total_documents: int
    total_chunks: int
    last_updated: str
    documents: List[Dict[str, str]]


class LMStudioEmbeddings:
    """LM Studio Embeddings 适配器"""
    
    def __init__(self, model: str = "text-embedding-qwen3-embedding-4b", base_url: str = "http://127.0.0.1:1234", api_key: str = None):
        self.model = model
        self.base_url = base_url
        self.api_key = api_key
        self.dimension = 2560  # Qwen3-Embedding-4B 输出维度
    
    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        embeddings = []
        for text in texts:
            result = self._embed(text)
            embeddings.append(result)
        return embeddings
    
    def embed_query(self, text: str) -> List[float]:
        return self._embed(text)
    
    def _embed(self, text: str) -> List[float]:
        try:
            payload = {
                "model": self.model,
                "input": text,
                "encoding_format": "float"
            }
            # 构建 headers，支持 API Key 认证
            headers = {}
            if self.api_key:
                headers["Authorization"] = f"Bearer {self.api_key}"

            response = requests.post(
                f"{self.base_url}/v1/embeddings",
                json=payload,
                headers=headers,
                timeout=60
            )
            if response.status_code == 200:
                result = response.json()
                embedding = result.get("data", [{}])[0].get("embedding", [])
                if not embedding:
                    raise Exception("LM Studio 返回空 Embedding")
                return embedding
            else:
                raise Exception(f"LM Studio Embedding 失败：{response.status_code}")
        except Exception as e:
            print(f"[Error] Embedding 错误：{e}")
            raise

    def embed_documents(self, texts: List[str], max_retries: int = 2) -> List[List[float]]:
        """批量嵌入文档，单条失败时自动重试"""
        embeddings = []
        for i, text in enumerate(texts):
            last_error = None
            for attempt in range(max_retries + 1):
                try:
                    result = self._embed(text)
                    embeddings.append(result)
                    break
                except Exception as e:
                    last_error = e
                    if attempt < max_retries:
                        print(f"[Warning] 第 {i+1}/{len(texts)} 条文本 Embedding 失败（尝试 {attempt+1}/{max_retries+1}）：{e}，正在重试...")
                    else:
                        raise Exception(f"第 {i+1}/{len(texts)} 条文本 Embedding 失败，已重试 {max_retries} 次：{last_error}") from last_error
        return embeddings


class KnowledgeBase:
    def __init__(
        self,
        name: str = "default",
        documents_path: Optional[str] = None,
        embed_model: str = "text-embedding-qwen3-embedding-4b",
        llm_model: str = "qwen/qwen3.5-9b",
        chunk_size: int = 800,
        chunk_overlap: int = 100,
        lmstudio_base_url: str = "http://127.0.0.1:1234",
        qdrant_host: str = "localhost",
        qdrant_port: int = 6333,
        collection_name: str = "kb_default",
        api_key: str = None
    ):
        self.name = name

        # 动态路径推导：core/ 的父目录为项目根目录
        if documents_path is None:
            CORE_DIR = os.path.dirname(os.path.abspath(__file__))
            INSTALL_DIR = os.path.dirname(CORE_DIR)
            documents_path = os.path.join(INSTALL_DIR, "documents")

        self.documents_path = documents_path
        self.embed_model = embed_model
        self.llm_model = llm_model
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.lmstudio_base_url = lmstudio_base_url
        self.api_key = api_key
        self.collection_name = collection_name
        
        # 初始化 Embedding
        self.embeddings = LMStudioEmbeddings(model=embed_model, base_url=lmstudio_base_url, api_key=api_key)
        
        # 初始化 Qdrant 客户端
        self.client = QdrantClient(host=qdrant_host, port=qdrant_port)
        
        # 创建或获取 Collection
        self._init_collection()
        
        # 文档索引缓存
        self._doc_cache = {}
    
    def _init_collection(self):
        """初始化 Qdrant Collection"""
        try:
            # 检查 Collection 是否存在
            collections = self.client.get_collections().collections
            collection_names = [c.name for c in collections]
            
            if self.collection_name not in collection_names:
                # 创建 Collection
                self.client.create_collection(
                    collection_name=self.collection_name,
                    vectors_config=VectorParams(
                        size=self.embeddings.dimension,  # 2560 维
                        distance=Distance.COSINE
                    )
                )
                print(f"[INFO] 创建 Collection: {self.collection_name}")
            else:
                print(f"[INFO] Collection 已存在：{self.collection_name}")
        except Exception as e:
            print(f"[ERROR] 初始化 Collection 失败：{e}")
            raise
    
    def _chunk_text(self, text: str) -> List[str]:
        """文本分块"""
        # 防止错误配置导致无限循环
        if self.chunk_overlap >= self.chunk_size:
            print(f"[WARN] chunk_overlap({self.chunk_overlap}) >= chunk_size({self.chunk_size})，自动调整 chunk_overlap")
            self.chunk_overlap = max(0, self.chunk_size // 4)

        chunks = []
        start = 0
        text_length = len(text)

        while start < text_length:
            end = min(start + self.chunk_size, text_length)
            chunk = text[start:end]

            # 如果不是最后一块，尝试在句子边界处分割
            if end < text_length:
                # 寻找最近的句子结束符
                for sep in ['。', '！', '？', '!', '?', '.', '\n']:
                    last_sep = chunk.rfind(sep)
                    if last_sep > self.chunk_size // 2:
                        chunk = chunk[:last_sep + 1]
                        break

            chunks.append(chunk.strip())
            # 确保每次至少前进 1 个字符
            next_start = end - self.chunk_overlap if end < text_length else text_length
            start = max(next_start, start + 1)

        return chunks
    
    def _load_pdf(self, filepath: str) -> str:
        """读取 PDF 文件（优先读取 OCR 文本）"""
        # 优先检查 OCR 文本文件
        ocr_txt_path = filepath.replace('.pdf', '_ocr.txt')
        if os.path.exists(ocr_txt_path):
            with open(ocr_txt_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        # 否则直接读取 PDF 文本
        text = ""
        try:
            with pymupdf.open(filepath) as doc:
                for page in doc:
                    text += page.get_text()
        except Exception as e:
            print(f"[Warning] 读取 PDF 失败 {filepath}: {e}")
        return text
    
    def _load_docx(self, filepath: str) -> str:
        """读取 Word 文档"""
        try:
            from docx import Document
            doc = Document(filepath)
            text = ""
            
            # 读取段落
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # 读取表格
            for table in doc.tables:
                for row in table.rows:
                    row_text = ""
                    for cell in row.cells:
                        row_text += cell.text + "\t"
                    text += row_text + "\n"
                text += "\n"
            
            return text
        except Exception as e:
            print(f"[Warning] 读取 Word 失败 {filepath}: {e}")
            return ""
    
    def _load_txt(self, filepath: str) -> str:
        """读取 TXT 文件"""
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            print(f"[Warning] 读取 TXT 失败 {filepath}: {e}")
            return ""
    
    def _load_excel(self, filepath: str) -> str:
        """读取 Excel 文件"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath)
            text = ""
            
            for sheet in wb.worksheets:
                text += f"=== {sheet.title} ===\n"
                for row in sheet.iter_rows(values_only=True):
                    text += "\t".join(str(cell) if cell is not None else "" for cell in row) + "\n"
                text += "\n"
            
            return text
        except Exception as e:
            print(f"[Warning] 读取 Excel 失败 {filepath}: {e}")
            return ""
    
    def _load_pptx(self, filepath: str) -> str:
        """读取 PPT 文件"""
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text = ""
            
            for i, slide in enumerate(prs.slides, 1):
                text += f"=== 第{i}页 ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            
            return text
        except Exception as e:
            print(f"[Warning] 读取 PPT 失败 {filepath}: {e}")
            return ""
    
    def _load_document(self, filepath: str) -> str:
        """根据文件类型加载文档"""
        ext = os.path.splitext(filepath)[1].lower()
        
        loaders = {
            '.pdf': self._load_pdf,
            '.docx': self._load_docx,
            '.txt': self._load_txt,
            '.xlsx': self._load_excel,
            '.xlsm': self._load_excel,
            '.pptx': self._load_pptx,
        }
        
        loader = loaders.get(ext)
        if loader:
            return loader(filepath)
        else:
            print(f"[Warning] 不支持的文件类型：{ext}")
            return ""
    
    def _generate_doc_id(self, filepath: str) -> str:
        """生成文档 ID（字符串，用于缓存和追踪）"""
        return hashlib.md5(filepath.encode()).hexdigest()
    
    def _generate_point_id(self, doc_id: str, chunk_idx: int) -> str:
        """生成 Qdrant 点 ID（字符串 UUID 形式，避免整数溢出和冲突）"""
        # 使用 doc_id + chunk_idx 的 MD5 哈希作为唯一 ID
        return hashlib.md5(f"{doc_id}_{chunk_idx}".encode()).hexdigest()
    
    def refresh_index(self, force: bool = False, auto_ocr: bool = True):
        """刷新索引"""
        print("\n" + "="*60)
        print("开始刷新索引...")
        print("="*60)
        
        start_time = time.time()
        
        # 扫描文档
        patterns = ['*.pdf', '*.docx', '*.txt', '*.xlsx', '*.xlsm', '*.pptx']
        documents = []
        
        for pattern in patterns:
            files = glob.glob(os.path.join(self.documents_path, '**', pattern), recursive=True)
            # 排除 OCR 生成的 txt 文件
            files = [f for f in files if not f.endswith('_ocr.txt')]
            documents.extend(files)
        
        print(f"[INFO] 扫描到 {len(documents)} 个文档")
        
        # 处理每个文档
        total_chunks = 0
        for filepath in documents:
            try:
                doc_id = self._generate_doc_id(filepath)
                filename = os.path.basename(filepath)
                
                # 如果非强制刷新且文档已存在，跳过
                if not force and doc_id in self._doc_cache:
                    print(f"[SKIP] {filename}")
                    continue
                
                # 加载文档
                print(f"[LOAD] {filename}")
                content = self._load_document(filepath)
                
                if not content:
                    print(f"[WARN] 文档内容为空：{filename}")
                    continue
                
                # 分块
                chunks = self._chunk_text(content)
                print(f"[CHUNK] {filename}: {len(chunks)} 块")
                
                # 生成向量
                embeddings = self.embeddings.embed_documents(chunks)
                
                # 准备 Qdrant Points
                points = []
                for idx, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                    point_id = self._generate_point_id(doc_id, idx)
                    point = PointStruct(
                        id=point_id,
                        vector=embedding,
                        payload={
                            "text": chunk,
                            "source": filename,
                            "filepath": filepath,
                            "doc_id": doc_id,
                            "chunk_idx": idx
                        }
                    )
                    points.append(point)
                
                # 删除旧数据
                try:
                    self.client.delete(
                        collection_name=self.collection_name,
                        points_selector=Filter(
                            must=[FieldCondition(key="doc_id", match=MatchValue(value=doc_id))]
                        )
                    )
                except:
                    pass
                
                # 上传到 Qdrant
                if points:
                    self.client.upsert(
                        collection_name=self.collection_name,
                        points=points
                    )
                    total_chunks += len(points)
                    print(f"[OK] {filename}: {len(points)} 块")
                
                # 更新缓存
                self._doc_cache[doc_id] = {
                    "filename": filename,
                    "filepath": filepath,
                    "chunks": len(chunks),
                    "time": datetime.now().isoformat()
                }
                
            except Exception as e:
                print(f"[ERROR] 处理文档失败 {filepath}: {e}")
        
        elapsed = time.time() - start_time
        
        print("\n" + "="*60)
        print(f"索引刷新完成！")
        print(f"  文档数：{len(self._doc_cache)}")
        print(f"  文本块数：{total_chunks}")
        print(f"  耗时：{elapsed:.2f} 秒")
        print("="*60 + "\n")
        
        return {
            "total_documents": len(self._doc_cache),
            "total_chunks": total_chunks,
            "elapsed": elapsed
        }
    
    def query(self, question: str, top_k: int = 2, auto_refresh: bool = False) -> QueryResult:
        """查询知识库"""
        start_time = time.time()
        
        # 自动刷新（可选）
        if auto_refresh:
            self.refresh_index()
        
        # 生成查询向量
        query_embedding = self.embeddings.embed_query(question)
        
        # Qdrant 查询
        search_results = self.client.query_points(
            collection_name=self.collection_name,
            query=query_embedding,
            limit=top_k
        )
        
        # 整理结果
        sources = []
        context_texts = []
        
        for result in search_results.points:
            payload = result.payload
            sources.append({
                "file": payload.get("source", "Unknown"),
                "content": payload.get("text", ""),
                "score": float(result.score) if hasattr(result, 'score') else 0.0
            })
            context_texts.append(f"[来源：{payload.get('source', 'Unknown')}]\n{payload.get('text', '')}")
        
        # 构建上下文
        context = "\n\n".join(context_texts)
        
        # 调用 LLM 生成答案
        answer, prompt_tokens, completion_tokens = self._call_llm(question, context)
        
        elapsed = time.time() - start_time
        
        return QueryResult(
            answer=answer,
            sources=sources,
            query=question,
            timestamp=datetime.now().isoformat(),
            prompt_tokens=prompt_tokens,
            completion_tokens=completion_tokens,
            total_tokens=prompt_tokens + completion_tokens,
            query_time=elapsed
        )
    
    def _call_llm(self, question: str, context: str) -> tuple:
        """调用 LLM 生成答案"""
        prompt = f"""基于以下上下文信息回答问题。如果上下文中没有相关信息，请说明你不知道。

上下文信息：
{context}

问题：{question}

请用中文回答，并标注信息来源。

答案："""
        
        try:
            # 构建 headers，支持 API Key 认证
            headers = {}
            if self.api_key:
                headers["Authorization"] = f"Bearer {self.api_key}"
            
            response = requests.post(
                f"{self.lmstudio_base_url}/v1/completions",
                json={
                    "model": self.llm_model,
                    "prompt": prompt,
                    "max_tokens": 1024,
                    "temperature": 0.7,
                    "stream": False
                },
                headers=headers,
                timeout=60
            )
            
            if response.status_code == 200:
                result = response.json()
                answer = result.get("choices", [{}])[0].get("text", "").strip()
                usage = result.get("usage", {})
                prompt_tokens = usage.get("prompt_tokens", 0)
                completion_tokens = usage.get("completion_tokens", 0)
                return answer, prompt_tokens, completion_tokens
            else:
                return f"[LLM 调用失败：{response.status_code}]", 0, 0
                
        except Exception as e:
            return f"[LLM 调用错误：{e}]", 0, 0
    
    def get_status(self) -> IndexStatus:
        """获取索引状态"""
        try:
            # 从 Qdrant 获取总文本块数
            count = self.client.count(collection_name=self.collection_name)
            total_chunks = count.count
            
            # 简化处理：返回缓存的文档信息
            documents = []
            for doc_id, info in self._doc_cache.items():
                documents.append({
                    "filename": info.get("filename", "Unknown"),
                    "doc_id": doc_id,
                    "chunks": info.get("chunks", 0)
                })
            
            # 如果缓存为空但实际有数据，尝试获取一次
            if len(documents) == 0 and total_chunks > 0:
                # 简单查询获取一个样本
                sample = self.client.query_points(
                    collection_name=self.collection_name,
                    query=[0.0] * 2560,
                    limit=1,
                    with_payload=["source", "doc_id"]
                )
                if sample.points:
                    p = sample.points[0]
                    documents.append({
                        "filename": p.payload.get("source", "Unknown"),
                        "doc_id": p.payload.get("doc_id", ""),
                        "chunks": total_chunks
                    })
            
            return IndexStatus(
                total_documents=len(documents),
                total_chunks=total_chunks,
                last_updated=datetime.now().isoformat(),
                documents=documents
            )
        except Exception as e:
            print(f"[ERROR] 获取状态失败：{e}")
            import traceback
            traceback.print_exc()
            return IndexStatus(
                total_documents=0,
                total_chunks=0,
                last_updated="",
                documents=[]
            )
    
    def delete_collection(self):
        """删除 Collection（用于重建）"""
        try:
            self.client.delete_collection(collection_name=self.collection_name)
            print(f"[INFO] 已删除 Collection: {self.collection_name}")
        except Exception as e:
            print(f"[ERROR] 删除 Collection 失败：{e}")


# 便捷函数
def create_knowledge_base(**kwargs) -> KnowledgeBase:
    """创建知识库实例"""
    return KnowledgeBase(**kwargs)


def query_knowledge_base(question: str, **kwargs) -> QueryResult:
    """快速查询"""
    kb = create_knowledge_base()
    return kb.query(question, **kwargs)


def refresh_index(**kwargs):
    """刷新索引"""
    kb = create_knowledge_base()
    return kb.refresh_index(**kwargs)


def get_status() -> IndexStatus:
    """获取状态"""
    kb = create_knowledge_base()
    return kb.get_status()


if __name__ == "__main__":
    # 测试
    print("="*60)
    print("RAG Core - Qdrant 版本测试")
    print("="*60)
    
    # 创建知识库
    kb = KnowledgeBase()
    
    # 刷新索引
    kb.refresh_index(force=True)
    
    # 查询测试
    result = kb.query("知识库中有什么文档？", top_k=2)
    print(f"\n问题：{result.query}")
    print(f"答案：{result.answer}")
    print(f"耗时：{result.query_time:.2f}秒")
    print(f"Token: {result.total_tokens}")
