"""
RAG Core - 简化版 (LM Studio 后端，最小依赖)
"""
import os
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

import glob
import hashlib
import json
from datetime import datetime
from typing import List, Dict, Optional, Any, Tuple
from dataclasses import dataclass
import requests
import chromadb
from chromadb.config import Settings

import pymupdf
import time

from .excel_analyzer import ExcelAnalyzer, format_excel_result

@dataclass
class QueryResult:
    answer: str
    sources: List[Dict[str, Any]]
    query: str
    timestamp: str
    # Token 统计和耗时
    prompt_tokens: int = 0
    completion_tokens: int = 0
    total_tokens: int = 0
    query_time: float = 0.0

@dataclass
class IndexStatus:
    total_documents: int
    total_chunks: int
    last_updated: str
    documents: List[Dict[str, str]]


class LMStudioEmbeddings:
    """LM Studio Embeddings 适配器（兼容 ChromaDB）"""
    
    def __init__(self, model: str = "text-embedding-qwen3-embedding-4b", base_url: str = "http://127.0.0.1:1234", api_key: str = None):
        self.model = model
        self.base_url = base_url
        self.api_key = api_key
    
    def name(self) -> str:
        """ChromaDB 需要的 name 方法"""
        return f"lmstudio_{self.model}"
    
    def __call__(self, input: List[str]) -> List[List[float]]:
        """ChromaDB 调用接口"""
        return self.embed_documents(input)
    
    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        embeddings = []
        for text in texts:
            result = self._embed(text)
            embeddings.append(result)
        return embeddings
    
    def embed_query(self, text: str) -> List[float]:
        return self._embed(text)
    
    def _embed(self, text: str) -> List[float]:
        try:
            payload = {
                "model": self.model,
                "input": text,
                "encoding_format": "float"
            }
            # 构建 headers，支持 API Key 认证
            headers = {}
            if self.api_key:
                headers["Authorization"] = f"Bearer {self.api_key}"
            
            response = requests.post(
                f"{self.base_url}/v1/embeddings",
                json=payload,
                headers=headers,
                timeout=60
            )
            if response.status_code == 200:
                result = response.json()
                embedding = result.get("data", [{}])[0].get("embedding", [])
                if not embedding:
                    raise Exception("LM Studio 返回空 Embedding")
                return embedding
            else:
                raise Exception(f"LM Studio Embedding 失败：{response.status_code}")
        except Exception as e:
            print(f"[Error] Embedding 错误：{e}")
            raise


class KnowledgeBase:
    def __init__(
        self,
        name: str = "default",
        documents_path: Optional[str] = None,  # 动态计算默认路径
        db_path: Optional[str] = None,
        embed_model: Optional[str] = None,
        llm_model: Optional[str] = None,
        chunk_size: int = 800,
        chunk_overlap: int = 100,
        lmstudio_base_url: Optional[str] = None,
        embedding_base_url: Optional[str] = None,
        api_key: Optional[str] = None,
        embedding_api_key: Optional[str] = None
    ):
        self.name = name
        
        # ==================== 从 .env 读取配置（如果未提供） ====================
        CORE_DIR = os.path.dirname(os.path.abspath(__file__))
        INSTALL_DIR = os.path.dirname(CORE_DIR)
        env_path = os.path.join(INSTALL_DIR, "config", ".env")
        
        if os.path.exists(env_path):
            from dotenv import load_dotenv
            load_dotenv(env_path)
        
        # 使用传入值或环境变量
        self.api_key = api_key or os.getenv("API_KEY", None)
        self.lmstudio_base_url = lmstudio_base_url or os.getenv("LLM_BASE_URL", "http://127.0.0.1:1234")
        self.llm_model = llm_model or os.getenv("LLM_MODEL", "qwen/qwen3.5-9b")
        self.embed_model = embed_model or os.getenv("EMBEDDING_MODEL", "text-embedding-qwen3-embedding-4b")
        self.embedding_base_url = embedding_base_url or os.getenv("EMBEDDING_BASE_URL", self.lmstudio_base_url)
        self.embedding_api_key = embedding_api_key or os.getenv("EMBEDDING_API_KEY", self.api_key)
        
        # ==================== 动态路径计算 ====================
        # 如果 documents_path 为 None，自动计算
        if documents_path is None:
            # documents 目录
            documents_path = os.path.join(INSTALL_DIR, "documents")
        
        self.documents_path = documents_path
        if db_path is None:
            db_path = os.path.join(self.documents_path, ".rag", "chroma_db")
        self.db_path = db_path
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
        self.meta_path = os.path.join(documents_path, ".rag", "metadata.json")
        
        # 初始化 Embeddings（使用从 .env 读取的配置）
        self.embeddings = LMStudioEmbeddings(model=self.embed_model, base_url=self.embedding_base_url, api_key=self.embedding_api_key)
        
        # 初始化 ChromaDB（使用自定义 Embedding 函数）
        os.makedirs(self.db_path, exist_ok=True)
        self.client = chromadb.PersistentClient(path=self.db_path)
        self.collection = self.client.get_or_create_collection(
            name=f"kb_{name}",
            embedding_function=self.embeddings
        )
        
        # 文本分块器
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
        # 初始化 OCR 引擎（用于 Word 图片识别）
        self.ocr = None

        # 初始化 Excel 结构化分析器
        self.excel_analyzer = ExcelAnalyzer(
            llm_base_url=self.lmstudio_base_url,
            llm_model=self.llm_model,
            api_key=self.api_key
        )

        print(f"[OK] 知识库 '{name}' 初始化完成")
        print(f"   - LLM: {self.llm_model}")
        print(f"   - Embedding: {self.embed_model}")
        print(f"   - LLM Server: {self.lmstudio_base_url}")
        if self.embedding_base_url != self.lmstudio_base_url:
            print(f"   - Embedding Server: {self.embedding_base_url}")
    
    def _get_file_hash(self, filepath: str) -> str:
        hasher = hashlib.md5()
        with open(filepath, 'rb') as f:
            hasher.update(f.read())
        return hasher.hexdigest()
    
    def _load_metadata(self) -> Dict:
        if os.path.exists(self.meta_path):
            with open(self.meta_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {"documents": {}, "last_updated": None}
    
    def _save_metadata(self, metadata: Dict):
        os.makedirs(os.path.dirname(self.meta_path), exist_ok=True)
        with open(self.meta_path, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)
    
    def _load_pdf(self, filepath: str) -> str:
        """读取 PDF 文件（优先读取 OCR 文本）"""
        # 优先检查 OCR 文本文件
        ocr_txt_path = filepath.replace('.pdf', '_ocr.txt')
        if os.path.exists(ocr_txt_path):
            print(f"  [OCR] 使用 OCR 文本：{os.path.basename(ocr_txt_path)}")
            with open(ocr_txt_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        # 否则直接读取 PDF 文本
        text = ""
        with pymupdf.open(filepath) as doc:
            for page in doc:
                text += page.get_text()
        return text
    
    def _load_docx(self, filepath: str) -> str:
        """读取 Word 文件（包含段落、表格和图片 OCR）"""
        try:
            from docx import Document
            from PIL import Image
            import io
            import numpy as np
            
            doc = Document(filepath)
            text = ""
            
            # 读取所有段落
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
            
            # 读取所有表格
            for i, table in enumerate(doc.tables, 1):
                text += f"\n=== 表格 {i} ===\n"
                for j, row in enumerate(table.rows):
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text += f"行{j}: {' | '.join(row_text)}\n"
                text += "\n"
            
            # 读取并 OCR 图片
            image_count = 0
            ocr_text = ""
            
            # 遍历文档中的所有关系（包括图片）
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        # 获取图片数据
                        image_blob = rel.target_part.blob
                        image = Image.open(io.BytesIO(image_blob))
                        
                        # 转换为 numpy 数组供 OCR 使用
                        image_array = np.array(image)
                        
                        # 初始化 OCR 引擎（如果还未初始化）
                        if self.ocr is None:
                            print(f"  [OCR] 初始化 OCR 引擎（首次使用）...")
                            from rapidocr import RapidOCR
                            import rapidocr
                            # 使用 wheel 包自带的 infer 模型（无需联网下载）
                            models_dir = str(os.path.join(os.path.dirname(rapidocr.__file__), 'models'))
                            params = {
                                'Global.model_root_dir': models_dir,  # 必须显式设置，避免 WindowsPath 类型
                                'Det.model_path': str(os.path.join(models_dir, 'ch_PP-OCRv4_det_infer.onnx')),
                                'Rec.model_path': str(os.path.join(models_dir, 'ch_PP-OCRv4_rec_infer.onnx')),
                                'Cls.model_path': str(os.path.join(models_dir, 'ch_ppocr_mobile_v2.0_cls_infer.onnx')),
                            }
                            self.ocr = RapidOCR(params=params)
                            print(f"  [OCR] OCR 引擎初始化完成")
                        
                        # OCR 识别
                        print(f"  [OCR] 识别图片 {image_count + 1}...")
                        ocr_output = self.ocr(image_array)
                        
                        if ocr_output and ocr_output.txts:
                            ocr_text += f"\n[图片 {image_count + 1} 内容]\n"
                            for txt in ocr_output.txts:
                                if txt:
                                    ocr_text += txt + "\n"
                            ocr_text += "[/图片内容]\n"
                            print(f"  [OCR] 图片 {image_count + 1} 识别完成")
                        else:
                            print(f"  [OCR] 图片 {image_count + 1} 未识别出文字")
                        
                        image_count += 1
                    except Exception as e:
                        print(f"  [Warning] OCR 识别图片失败：{e}")
                        # 继续处理其他图片，不中断整个文档读取
                        continue
            
            if image_count > 0:
                text += f"\n\n=== 图片 OCR 内容（共 {image_count} 张）===\n"
                text += ocr_text
            
            # 即使没有 OCR 内容，也要返回文字部分
            if not text.strip():
                print(f"  [WARN] Word 文档无文字内容")
            
            return text
        except Exception as e:
            print(f"[Error] 读取 Word 文件失败：{e}")
            import traceback
            traceback.print_exc()
            return ""
    
    def _load_excel(self, filepath: str) -> str:
        """读取 Excel 文件"""
        try:
            import pandas as pd
            # 尝试读取所有 sheet
            excel_file = pd.ExcelFile(filepath)
            text = ""
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(filepath, sheet_name=sheet_name)
                text += f"\n=== {sheet_name} ===\n"
                text += df.to_string() + "\n"
            return text
        except Exception as e:
            print(f"[Warning] 读取 Excel 文件失败：{e}")
            return ""
    
    def _load_csv(self, filepath: str) -> str:
        """读取 CSV 文件"""
        try:
            import pandas as pd
            df = pd.read_csv(filepath)
            return df.to_string()
        except Exception as e:
            print(f"[Warning] 读取 CSV 文件失败：{e}")
            return ""
    
    def _load_pptx(self, filepath: str) -> str:
        """读取 PowerPoint 文件"""
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text = ""
            for i, slide in enumerate(prs.slides, 1):
                text += f"\n=== 幻灯片 {i} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"[Warning] 读取 PPT 文件失败：{e}")
            return ""
    
    def _load_txt(self, filepath: str) -> str:
        """读取 TXT/Markdown 文件"""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            print(f"[Warning] 读取文本文件失败：{e}")
            return ""
    
    def _load_document(self, filepath: str) -> str:
        """根据文件类型加载文档"""
        lower_path = filepath.lower()
        
        if lower_path.endswith(".pdf"):
            return self._load_pdf(filepath)
        elif lower_path.endswith(".docx") or lower_path.endswith(".doc"):
            return self._load_docx(filepath)
        elif lower_path.endswith(".xlsx") or lower_path.endswith(".xls"):
            return self._load_excel(filepath)
        elif lower_path.endswith(".csv"):
            return self._load_csv(filepath)
        elif lower_path.endswith(".pptx") or lower_path.endswith(".ppt"):
            return self._load_pptx(filepath)
        elif lower_path.endswith((".txt", ".md", ".markdown")):
            return self._load_txt(filepath)
        else:
            # 其他格式尝试作为文本读取
            return self._load_txt(filepath)
    
    def _split_text(self, text: str) -> List[str]:
        """简单文本分块"""
        # 防止错误配置导致无限循环
        if self.chunk_overlap >= self.chunk_size:
            print(f"[WARN] chunk_overlap({self.chunk_overlap}) >= chunk_size({self.chunk_size})，自动调整 chunk_overlap")
            self.chunk_overlap = max(0, self.chunk_size // 4)

        chunks = []
        start = 0
        text_len = len(text)
        while start < text_len:
            end = min(start + self.chunk_size, text_len)
            chunk = text[start:end]
            # 尝试在句子边界切断
            if end < text_len:
                for sep in ["\n\n", "\n", "。", "！", "？"]:
                    pos = chunk.rfind(sep)
                    if pos > self.chunk_size // 2:
                        chunk = chunk[:pos + len(sep)]
                        break
            chunks.append(chunk.strip())
            # 确保每次至少前进 1 个字符，避免 chunk_overlap >= chunk_size 时死循环
            next_start = end - self.chunk_overlap if end < text_len else text_len
            start = max(next_start, start + 1)
        return chunks
    
    def _scan_documents(self) -> List[str]:
        """扫描文档目录 - 支持所有常见格式（包括 OCR 文本）"""
        documents = []
        
        # 定义支持的格式
        supported_extensions = [
            '.pdf', '.PDF',
            '.docx', '.DOCX', '.doc', '.DOC',
            '.txt', '.TXT', '.md', '.MD', '.markdown', '.MARKDOWN',
            '.xlsx', '.XLSX', '.xls', '.XLS', '.csv', '.CSV',
            '.pptx', '.PPTX', '.ppt', '.PPT'
        ]
        
        # 递归扫描所有子目录
        for root, dirs, files in os.walk(self.documents_path):
            # 跳过隐藏目录和.rag 目录
            dirs[:] = [d for d in dirs if not d.startswith('.')]
            
            for file in files:
                filepath = os.path.join(root, file)
                filename = file.lower()
                ext = os.path.splitext(file)[1]
                
                # 优先使用 OCR 文本文件（如果有）
                if ext.lower() == '.pdf':
                    # 检查是否有对应的 OCR 文本文件（使用与扩展名无关的路径）
                    ocr_txt = self._get_ocr_txt_path(filepath)
                    if os.path.exists(ocr_txt):
                        documents.append(ocr_txt)  # 使用 OCR 文本
                        continue
                
                if ext in supported_extensions:
                    documents.append(filepath)
        
        # 去重并返回
        unique_docs = list(set(documents))
        
        if unique_docs:
            print(f"\n[INFO] 扫描到 {len(unique_docs)} 个文档:")
            for doc in unique_docs:
                print(f"  - {os.path.basename(doc)}")
        
        return unique_docs
    
    def _detect_need_ocr(self, filepath: str) -> bool:
        """检测文件是否需要 OCR 处理"""
        ext = os.path.splitext(filepath)[1].lower()
        
        # PDF 文件检测
        if ext == '.pdf':
            try:
                doc = pymupdf.open(filepath)
                total_text = ""
                # 检查前 3 页（如果页数不足则检查全部），避免封面导致的误判
                pages_to_check = min(3, len(doc))
                for page_idx in range(pages_to_check):
                    total_text += doc[page_idx].get_text()
                doc.close()

                # 如果检查页的文字总量很少，判定为图片 PDF
                if not total_text.strip() or len(total_text.strip()) < 100:
                    return True
            except Exception as e:
                print(f"  [WARN] 检测 PDF 失败：{e}")
                return False
        
        # Word 文档图片检测
        if ext in ['.docx', '.doc']:
            try:
                from docx import Document
                doc = Document(filepath)
                # 检查是否有图片
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        return True  # 有图片，需要 OCR
            except Exception as e:
                print(f"  [WARN] 检测 Word 文档图片失败：{e}")
                return False
        
        # PPT 文档图片检测
        if ext in ['.pptx', '.ppt']:
            try:
                from pptx import Presentation
                from pptx.oxml.xml_helper import parse_xml_string
                from pptx.oxml.ns import nsmap
                
                prs = Presentation(filepath)
                # 遍历所有幻灯片
                for slide in prs.slides:
                    # 检查是否有图片
                    for shape in slide.shapes:
                        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            return True
            except Exception as e:
                print(f"  [WARN] 检测 PPT 文档图片失败：{e}")
                return False
        
        return False
    
    def _get_ocr_txt_path(self, filepath: str) -> str:
        """根据原文件路径生成对应的 OCR 文本文件路径（不区分大小写扩展名）"""
        base, _ = os.path.splitext(filepath)
        return base + '_ocr.txt'

    def _auto_ocr(self, filepath: str) -> str:
        """自动对文件进行 OCR 处理

        Returns:
            OCR 文本文件路径，如果不需要 OCR 则返回原路径
        """
        ext = os.path.splitext(filepath)[1].lower()

        # 检查是否已有 OCR 文件
        ocr_txt = self._get_ocr_txt_path(filepath)
        if os.path.exists(ocr_txt):
            print(f"  [OCR] 已有 OCR 文件：{os.path.basename(ocr_txt)}")
            return ocr_txt

        # 检测是否需要 OCR
        if self._detect_need_ocr(filepath):
            print(f"  [OCR] 检测到图片文档，开始自动识别：{os.path.basename(filepath)}")

            # PDF 文件使用 ocr_simple 处理
            if ext == '.pdf':
                try:
                    import sys
                    # 正确添加 tools 目录到 sys.path
                    CORE_DIR = os.path.dirname(os.path.abspath(__file__))
                    INSTALL_DIR = os.path.dirname(CORE_DIR)
                    TOOLS_DIR = os.path.join(INSTALL_DIR, "tools")
                    if TOOLS_DIR not in sys.path:
                        sys.path.insert(0, TOOLS_DIR)
                    from ocr_simple import ocr_pdf

                    ocr_output = ocr_pdf(filepath)

                    if ocr_output:
                        print(f"  [OCR] ✅ 识别完成：{os.path.basename(ocr_output)}")
                        return ocr_output
                    else:
                        print(f"  [OCR] ⚠️ 识别失败，使用原文档")
                        return filepath

                except Exception as e:
                    print(f"  [OCR] ⚠️ OCR 处理失败：{e}，使用原文档")
                    return filepath

            # Word 和 PPT 文档不需要生成独立的 OCR 文件
            # OCR 处理在 _load_docx 和 _load_pptx 中进行
            # 这里只返回原路径，表示需要 OCR 但已标记
            print(f"  [OCR] ℹ️  文档包含图片，将在读取时自动 OCR")
            return filepath

        return filepath
    
    def add_documents(self, force: bool = False, auto_ocr: bool = True):
        """增量添加文档（自动检测并处理 OCR）
        
        Args:
            force: 是否强制重建索引
            auto_ocr: 是否自动检测并处理需要 OCR 的文件
        """
        print("\n[INFO] 检查文档更新...")
        
        metadata = self._load_metadata()
        indexed = metadata.get("documents", {})
        files = self._scan_documents()
        
        # ==================== 清理已删除文件的记录 ====================
        to_delete = []
        for filepath in indexed:
            if not os.path.exists(filepath):
                to_delete.append(filepath)
        
        if to_delete:
            print(f"\n  [CLEAN] 发现 {len(to_delete)} 个已删除文件，正在清理...")
            
            # 获取所有 chunks 的 ID，按文件路径前缀匹配删除
            try:
                all_items = self.collection.get(include=['metadatas'])
                ids_to_delete = []
                
                for filepath in to_delete:
                    # 通过 source 字段匹配（相对路径更可靠）
                    source = os.path.relpath(filepath, self.documents_path)
                    file_ids = [
                        id for id, meta in zip(all_items['ids'], all_items['metadatas'])
                        if meta.get('source', '') == source or meta.get('filepath', '') == filepath
                    ]
                    if file_ids:
                        ids_to_delete.extend(file_ids)
                        print(f"    • {os.path.basename(filepath)} ({len(file_ids)} 个块)")
                    else:
                        print(f"    • {os.path.basename(filepath)} (无向量记录)")
                
                if ids_to_delete:
                    self.collection.delete(ids=ids_to_delete)
                    print(f"  [OK] 删除 {len(ids_to_delete)} 个向量")
            except Exception as e:
                print(f"  [WARN] 向量清理失败: {e}")
            
            # 从 metadata 中删除
            for filepath in to_delete:
                del indexed[filepath]
            
            metadata["documents"] = indexed
            metadata["last_updated"] = datetime.now().isoformat()
            self._save_metadata(metadata)
            print(f"  [OK] 已清理 {len(to_delete)} 个无效记录\n")
        
        # ==================== 检查需要索引的文件 ====================
        to_index = []
        ocr_processed = []
        
        for filepath in files:
            # 自动 OCR 处理
            if auto_ocr:
                processed_path = self._auto_ocr(filepath)
                if processed_path != filepath:
                    ocr_processed.append((filepath, processed_path))
                    filepath = processed_path
            
            file_hash = self._get_file_hash(filepath)
            if filepath not in indexed or indexed[filepath] != file_hash or force:
                to_index.append(filepath)
                print(f"  [NEW] {os.path.basename(filepath)}")
        
        if not to_index:
            print("  [OK] 文档已是最新")
            return
        
        # 加载并分块文档
        all_chunks = []
        all_ids = []
        all_metadatas = []
        
        for filepath in to_index:
            # 根据文件类型加载
            text = self._load_document(filepath)

            # 对 Excel/CSV 文件额外保存 schema，用于结构化分析
            ext = os.path.splitext(filepath)[1].lower()
            if ext in ['.xlsx', '.xls', '.csv']:
                try:
                    sheets = self.excel_analyzer.load(filepath)
                    self.excel_analyzer.save_schema(filepath, sheets)
                    print(f"  [Excel] 已保存 schema：{os.path.basename(filepath)}")
                except Exception as e:
                    print(f"  [Warning] 保存 Excel schema 失败：{e}")

            if not text:
                print(f"  [SKIP] {os.path.basename(filepath)} (无法读取)")
                continue

            chunks = self._split_text(text)
            for i, chunk in enumerate(chunks):
                chunk_id = f"{filepath}_{i}"
                all_chunks.append(chunk)
                all_ids.append(chunk_id)
                all_metadatas.append({
                    "source": os.path.relpath(filepath, self.documents_path),
                    "filepath": filepath,
                    "hash": self._get_file_hash(filepath)
                })
        
        # 批量嵌入并添加
        if all_chunks:
            print(f"  正在嵌入 {len(all_chunks)} 个文本块...")
            embeddings = self.embeddings.embed_documents(all_chunks)
            
            self.collection.add(
                documents=all_chunks,
                embeddings=embeddings,
                ids=all_ids,
                metadatas=all_metadatas
            )
            
            # 更新 metadata
            for filepath in to_index:
                indexed[filepath] = self._get_file_hash(filepath)
            metadata["documents"] = indexed
            metadata["last_updated"] = datetime.now().isoformat()
            self._save_metadata(metadata)
            
            # 输出 OCR 统计
            if ocr_processed:
                print(f"\n  [OCR] 本次自动处理 {len(ocr_processed)} 个文件:")
                for orig, ocr_file in ocr_processed:
                    print(f"    • {os.path.basename(orig)} → {os.path.basename(ocr_file)}")
            
            print(f"  [OK] 已索引 {len(to_index)} 个文件，{len(all_chunks)} 个文本块")
    
    def get_status(self) -> IndexStatus:
        """获取索引状态"""
        metadata = self._load_metadata()
        indexed_docs = metadata.get("documents", {})
        
        # 过滤已删除的文档（检查文件是否实际存在）
        actual_docs = []
        deleted_docs = []
        for filepath in indexed_docs.keys():
            if os.path.exists(filepath):
                actual_docs.append(filepath)
            else:
                deleted_docs.append(filepath)
        
        # 如果有已删除文档，清理 metadata 并同步 ChromaDB
        if deleted_docs:
            print(f"[INFO] 发现 {len(deleted_docs)} 个已删除文档，正在清理...")
            self._cleanup_deleted_documents(deleted_docs, indexed_docs, metadata)
        
        return IndexStatus(
            total_documents=len(actual_docs),
            total_chunks=self.collection.count(),
            last_updated=metadata.get("last_updated", "从未"),
            documents=actual_docs
        )
    
    def _cleanup_deleted_documents(self, deleted_docs: List[str], indexed_docs: Dict, metadata: Dict):
        """清理已删除文档的索引记录"""
        # 从 metadata 中移除
        for filepath in deleted_docs:
            if filepath in indexed_docs:
                del indexed_docs[filepath]
                print(f"  [CLEAN] 移除元数据记录：{os.path.basename(filepath)}")
        
        # 从 ChromaDB 中删除对应的 chunks
        try:
            # 获取所有 chunks 的 metadata
            all_items = self.collection.get()
            ids_to_delete = []
            for i, metadata_item in enumerate(all_items.get('metadatas', [])):
                source_path = metadata_item.get('filepath', '')
                if source_path in deleted_docs:
                    ids_to_delete.append(all_items['ids'][i])
            
            if ids_to_delete:
                self.collection.delete(ids=ids_to_delete)
                print(f"  [CLEAN] 删除 {len(ids_to_delete)} 个文本块")
            
            # 更新 metadata 文件
            metadata["documents"] = indexed_docs
            metadata["last_updated"] = datetime.now().isoformat()
            self._save_metadata(metadata)
            print(f"[OK] 清理完成，剩余 {len(indexed_docs)} 个有效文档")
        except Exception as e:
            print(f"[WARN] ChromaDB 清理失败：{e}")
    
    def delete_collection(self):
        """删除 Collection（用于初始化）"""
        try:
            # 删除 ChromaDB collection
            self.client.delete_collection(name=f"kb_{self.name}")
            print(f"[INFO] 已删除 Collection: kb_{self.name}")
            
            # 删除元数据文件
            if os.path.exists(self.meta_path):
                os.remove(self.meta_path)
                print(f"[INFO] 已删除元数据文件：{self.meta_path}")
            
            # 重新创建空的 collection
            self.collection = self.client.get_or_create_collection(
                name=f"kb_{self.name}",
                embedding_function=self.embeddings
            )
            print(f"[INFO] 已创建新的空 Collection")
        except Exception as e:
            print(f"[ERROR] 删除 Collection 失败：{e}")
            raise
    
    def search(self, query: str, top_k: int = 5):
        """搜索相关文档"""
        self.add_documents()
        query_embedding = self.embeddings.embed_query(query)
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k
        )
        
        # 转换为 Document 格式
        docs = []
        if results['documents'] and results['documents'][0]:
            for i, content in enumerate(results['documents'][0]):
                metadata = results['metadatas'][0][i] if results['metadatas'] else {}
                docs.append(type('Document', (), {'page_content': content, 'metadata': metadata})())
        return docs
    
    def _is_complex_query(self, question: str) -> bool:
        """
        判断是否为复杂查询
        复杂问题使用混合检索，简单问题使用向量检索
        """
        COMPLEX_KEYWORDS = [
            '为什么', '如何', '怎么', '比较', '分析', '关系',
            '区别', '影响', '原因', '过程', '步骤', '比', '对比', '差异',
            '哪些', '什么关系', '有何不同', '优缺点'
        ]
        COMPLEX_LENGTH_THRESHOLD = 15
        
        if len(question) > COMPLEX_LENGTH_THRESHOLD:
            return True
        
        if any(kw in question for kw in COMPLEX_KEYWORDS):
            return True
        
        if question.count('?') > 1 or question.count('？') > 1:
            return True
        
        return False
    
    def query(self, question: str, top_k: int = 5, auto_refresh: bool = False, max_context_length: int = 8000) -> QueryResult:
        """
        智能查询 - 自动判断检索策略

        Args:
            question: 用户问题
            top_k: 返回的文档块数量
            auto_refresh: 是否自动检查并刷新索引
            max_context_length: 最大上下文长度（token 数）

        Returns:
            QueryResult 对象
        """
        if self._is_complex_query(question):
            print(f"[INFO] 复杂问题，使用混合检索")
            return self.query_hybrid(question, top_k, auto_refresh, max_context_length=max_context_length)
        else:
            print(f"[INFO] 简单问题，使用向量检索")
            return self.query_with_sources(question, top_k, auto_refresh, max_context_length=max_context_length)
    
    def extract_data(self, question: str, fields: List[str] = None, top_k: int = 5, auto_refresh: bool = False) -> Dict[str, Any]:
        """
        从 Excel/CSV 文件中提取结构化数据

        Args:
            question: 用户问题，例如"2025 年销售额总和是多少？"
            fields: 期望提取的字段列表（可选，目前作为辅助提示）
            top_k: 检索相关文档数量
            auto_refresh: 是否自动刷新索引

        Returns:
            分析结果字典
        """
        if auto_refresh:
            self.add_documents()

        # 1. 检索相关文档
        query_embedding = self.embeddings.embed_query(question)
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k
        )

        if not results['documents'] or not results['documents'][0]:
            return {
                'error': '未找到相关文档',
                'answer': '未在知识库中找到相关 Excel/CSV 文件。'
            }

        # 2. 筛选出 Excel/CSV 文件路径
        excel_files = []
        seen = set()
        for i, metadata in enumerate(results['metadatas'][0]):
            filepath = metadata.get('filepath', '')
            ext = os.path.splitext(filepath)[1].lower()
            if ext in ['.xlsx', '.xls', '.csv'] and filepath not in seen:
                excel_files.append(filepath)
                seen.add(filepath)

        if not excel_files:
            return {
                'error': '未找到 Excel/CSV 文件',
                'answer': '检索到的文档不是 Excel/CSV 格式，无法执行结构化分析。'
            }

        # 3. 对第一个相关 Excel/CSV 文件执行分析
        target_file = excel_files[0]
        print(f"[Excel] 结构化分析文件：{os.path.basename(target_file)}")

        try:
            analysis_result = self.excel_analyzer.analyze(question, target_file)

            # 4. 如果指定了 fields，尝试从结果中提取对应字段
            if fields and 'data' in analysis_result:
                extracted = []
                for row in analysis_result['data']:
                    extracted_row = {}
                    for field in fields:
                        # 模糊匹配字段名
                        for key in row.keys():
                            if field.lower() in str(key).lower() or str(key).lower() in field.lower():
                                extracted_row[field] = row[key]
                                break
                    if extracted_row:
                        extracted.append(extracted_row)
                analysis_result['extracted_fields'] = extracted

            # 5. 生成自然语言答案
            answer = format_excel_result(analysis_result)
            analysis_result['answer'] = answer

            return analysis_result

        except Exception as e:
            import traceback
            traceback.print_exc()
            return {
                'error': f'分析失败：{str(e)}',
                'answer': f'Excel 结构化分析失败：{str(e)}'
            }

    def query_hybrid(self, question: str, top_k: int = 5, auto_refresh: bool = False, alpha: float = 0.5, max_context_length: int = 8000) -> QueryResult:
        """
        混合检索查询 - 向量 + 关键词

        Args:
            question: 用户问题
            top_k: 返回的文档块数量
            auto_refresh: 是否自动检查并刷新索引
            alpha: 向量检索权重 (0.5=各 50%)
            max_context_length: 最大上下文长度（token 数）
        """
        start_time = time.time()

        try:
            if auto_refresh:
                self.add_documents()

            # 动态导入 hybrid_search（确保 core 目录在 sys.path 中）
            CORE_DIR = os.path.dirname(os.path.abspath(__file__))
            if CORE_DIR not in sys.path:
                sys.path.insert(0, CORE_DIR)

            import hybrid_search
            HybridSearch = hybrid_search.HybridSearch

            hs = HybridSearch(self.collection, self.embeddings, alpha=alpha)
            hs.build_bm25_index()
            results = hs.search(question, top_k=top_k)

            return self._build_query_result(question, results, start_time, max_context_length=max_context_length)
        except Exception as e:
            import traceback
            print(f"[ERROR] 混合检索失败: {e}")
            traceback.print_exc()
            query_time = time.time() - start_time
            return QueryResult(
                answer=f"混合检索出错：{str(e)}",
                sources=[],
                query=question,
                timestamp=datetime.now().isoformat(),
                query_time=query_time
            )
    
    def _estimate_tokens(self, text: str) -> int:
        """简单估算 token 数（中文/数字/符号按 1 token，英文按 4 字符 1 token）"""
        tokens = 0
        for char in text:
            if '一' <= char <= '鿿':
                tokens += 1
            elif char.isdigit() or char.isspace() or not char.isalnum():
                tokens += 1
            else:
                tokens += 0.25
        return int(tokens)

    def _build_context(self, documents: List[str], metadatas: List[Dict], question: str, max_context_length: int = 8000) -> Tuple[str, List[Dict]]:
        """
        根据 max_context_length 构建上下文，自动截断避免超过模型上下文窗口

        Returns:
            context: 构建好的上下文字符串
            sources: 来源列表
        """
        # 预留 prompt 模板、问题和答案的空间（约 300 tokens）
        reserved_tokens = self._estimate_tokens(question) + 300
        available_tokens = max(0, max_context_length - reserved_tokens)

        context_parts = []
        sources = []
        current_tokens = 0

        for i, content in enumerate(documents):
            preview = content[:600]
            part = f"[来源{i+1}] {preview}"
            part_tokens = self._estimate_tokens(part)

            # 如果添加该来源会超出上限，且已有内容，则截断
            if current_tokens + part_tokens > available_tokens and context_parts:
                print(f"[INFO] 上下文接近上限，已使用 {current_tokens} tokens，截断后续来源")
                break

            context_parts.append(part)
            current_tokens += part_tokens

            metadata = metadatas[i] if metadatas else {}
            sources.append({
                "index": i + 1,
                "source": metadata.get("source", "未知"),
                "content_preview": preview[:150]
            })

        context = "\n\n".join(context_parts)
        return context, sources

    def query_with_sources(self, question: str, top_k: int = 5, auto_refresh: bool = False, max_context_length: int = 8000) -> QueryResult:
        """带引用来源的查询

        Args:
            question: 用户问题
            top_k: 返回的文档块数量
            auto_refresh: 是否自动检查并刷新索引（默认 False）
            max_context_length: 最大上下文长度（token 数，默认 8000，避免超过模型上下文窗口）
        """
        start_time = time.time()  # ⏱️ 记录开始时间

        # 可选：自动刷新索引（仅在添加新文档后需要）
        if auto_refresh:
            self.add_documents()

        # 检索
        query_embedding = self.embeddings.embed_query(question)
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k
        )

        if not results['documents'] or not results['documents'][0]:
            query_time = time.time() - start_time
            return QueryResult(
                answer="未在知识库中找到相关信息。",
                sources=[],
                query=question,
                timestamp=datetime.now().isoformat(),
                query_time=query_time
            )

        # 构建上下文 - 根据 max_context_length 自动截断
        context, sources = self._build_context(
            results['documents'][0],
            results['metadatas'][0] if results['metadatas'] else [],
            question,
            max_context_length
        )

        # 调用 LM Studio 生成答案
        prompt_tokens = 0
        completion_tokens = 0
        total_tokens = 0

        try:
            llm_prompt = f"""基于以下文档内容回答问题：

{context}

问题：{question}

要求：
- 简洁明了，列出关键要点
- 如果文档中没有相关信息，直接说"未在知识库中找到相关信息"

答案："""

            payload = {
                "model": self.llm_model,
                "messages": [{"role": "user", "content": llm_prompt}],
                "temperature": 0,
                "max_tokens": 500,
                "stream": False
            }
            # 构建 headers，支持 API Key 认证
            headers = {}
            if self.api_key:
                headers["Authorization"] = f"Bearer {self.api_key}"
            
            response = requests.post(
                f"{self.lmstudio_base_url}/v1/chat/completions",
                json=payload,
                headers=headers,
                timeout=60
            )
            if response.status_code == 200:
                result = response.json()
                message = result.get("choices", [{}])[0].get("message", {})
                
                # 直接使用 content 作为答案
                answer = message.get("content", "").strip()
                
                # ✅ 提取 token 信息
                usage = result.get("usage", {})
                prompt_tokens = usage.get("prompt_tokens", 0)
                completion_tokens = usage.get("completion_tokens", 0)
                total_tokens = usage.get("total_tokens", 0)
                
                # 清理重复内容
                if answer.count('\n\n答案：') > 1:
                    parts = answer.split('\n\n答案：')
                    answer = parts[-1].strip()
                
                # 最终兜底
                if not answer:
                    answer = "根据检索到的文档，相关内容如下：\n\n" + context[:600]
            else:
                answer = f"LM Studio 调用失败：{response.status_code}"
        except Exception as e:
            answer = f"LM Studio 调用错误：{e}"
        
        # ⏱️ 计算耗时
        query_time = time.time() - start_time
        
        return QueryResult(
            answer=answer,
            sources=sources,
            query=question,
            timestamp=datetime.now().isoformat(),
            prompt_tokens=prompt_tokens,
            completion_tokens=completion_tokens,
            total_tokens=total_tokens,
            query_time=query_time
        )
    
    def _build_query_result(self, question: str, results: Dict, start_time: float, max_context_length: int = 8000) -> QueryResult:
        """
        构建查询结果（供 query_hybrid 和 query_with_sources 共用）

        Args:
            question: 用户问题
            results: 检索结果（包含 ids, documents, metadatas, scores）
            start_time: 开始时间
            max_context_length: 最大上下文长度（token 数）
        """
        query_time = time.time() - start_time

        if not results.get('documents'):
            return QueryResult(
                answer="未在知识库中找到相关信息。",
                sources=[],
                query=question,
                timestamp=datetime.now().isoformat(),
                query_time=query_time
            )

        # 构建上下文 - 根据 max_context_length 自动截断
        context, sources = self._build_context(
            results['documents'],
            results['metadatas'] if results.get('metadatas') else [],
            question,
            max_context_length
        )
        
        # 调用 LM Studio 生成答案
        prompt_tokens = 0
        completion_tokens = 0
        total_tokens = 0
        
        try:
            llm_prompt = f"""基于以下文档内容回答问题：

{context}

问题：{question}

要求：
- 简洁明了，列出关键要点
- 如果文档中没有相关信息，直接说"未在知识库中找到相关信息"

答案："""

            payload = {
                "model": self.llm_model,
                "messages": [{"role": "user", "content": llm_prompt}],
                "temperature": 0,
                "max_tokens": 500,
                "stream": False
            }
            headers = {}
            if self.api_key:
                headers["Authorization"] = f"Bearer {self.api_key}"
            
            response = requests.post(
                f"{self.lmstudio_base_url}/v1/chat/completions",
                json=payload,
                headers=headers,
                timeout=60
            )
            if response.status_code == 200:
                result = response.json()
                message = result.get("choices", [{}])[0].get("message", {})
                answer = message.get("content", "").strip()
                
                usage = result.get("usage", {})
                prompt_tokens = usage.get("prompt_tokens", 0)
                completion_tokens = usage.get("completion_tokens", 0)
                total_tokens = usage.get("total_tokens", 0)
                
                if answer.count('\n\n答案：') > 1:
                    parts = answer.split('\n\n答案：')
                    answer = parts[-1].strip()
                
                if not answer:
                    answer = "根据检索到的文档，相关内容如下：\n\n" + context[:600]
            else:
                answer = f"LM Studio 调用失败：{response.status_code}"
        except Exception as e:
            answer = f"LM Studio 调用错误：{e}"
        
        return QueryResult(
            answer=answer,
            sources=sources,
            query=question,
            timestamp=datetime.now().isoformat(),
            prompt_tokens=prompt_tokens,
            completion_tokens=completion_tokens,
            total_tokens=total_tokens,
            query_time=query_time
        )


if __name__ == "__main__":
    print("=" * 60)
    print("RAG 知识库工具测试 - LM Studio 简化版")
    print("=" * 60)
    
    kb = KnowledgeBase()
    status = kb.get_status()
    print(f"\n索引状态:")
    print(f"  文档数：{status.total_documents}")
    print(f"  文本块数：{status.total_chunks}")
    print(f"  最后更新：{status.last_updated}")
    
    print("\n" + "=" * 60)
    print("测试查询：2026 年 2 月电费单的计费电量是多少？")
    print("=" * 60)
    
    result = kb.query_with_sources("2026 年 2 月电费单的计费电量是多少？")
    print(f"\n{result.answer}")
    
    if result.sources:
        print("\n来源文档:")
        seen = set()
        for s in result.sources:
            if s['source'] not in seen:
                seen.add(s['source'])
                print(f"  - {s['source']}")
