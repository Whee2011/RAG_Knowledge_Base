"""
RAG Core - 简化版 (LM Studio 后端，最小依赖)

⚠️  注意：此文件为历史备用实现，项目主入口为 core/core.py。
      新开发请使用 core/core.py，此文件仅作为向后兼容保留。
"""
import os
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

import glob
import hashlib
import json
from datetime import datetime
from typing import List, Dict, Optional, Any
from dataclasses import dataclass
import requests
import chromadb
from chromadb.config import Settings

import pymupdf
import time

@dataclass
class QueryResult:
    answer: str
    sources: List[Dict[str, Any]]
    query: str
    timestamp: str
    # Token 统计和耗时
    prompt_tokens: int = 0
    completion_tokens: int = 0
    total_tokens: int = 0
    query_time: float = 0.0

@dataclass
class IndexStatus:
    total_documents: int
    total_chunks: int
    last_updated: str
    documents: List[Dict[str, str]]


class LMStudioEmbeddings:
    """LM Studio Embeddings 适配器（兼容 ChromaDB）"""
    
    def __init__(self, model: str = "text-embedding-qwen3-embedding-4b", base_url: str = "http://127.0.0.1:1234", api_key: str = None):
        self.model = model
        self.base_url = base_url
        self.api_key = api_key
    
    def name(self) -> str:
        """ChromaDB 需要的 name 方法"""
        return f"lmstudio_{self.model}"
    
    def __call__(self, input: List[str]) -> List[List[float]]:
        """ChromaDB 调用接口"""
        return self.embed_documents(input)
    
    def _embed_with_retry(self, text: str, max_retries: int = 2, label: str = "") -> List[float]:
        """带重试的单个文本嵌入"""
        last_error = None
        prefix = f"{label} " if label else ""
        for attempt in range(max_retries + 1):
            try:
                return self._embed(text)
            except Exception as e:
                last_error = e
                if attempt < max_retries:
                    print(f"[Warning] {prefix}Embedding 失败（尝试 {attempt+1}/{max_retries+1}）：{e}，正在重试...")
                else:
                    raise last_error
        raise last_error

    def embed_documents(self, texts: List[str], max_retries: int = 2) -> List[List[float]]:
        """批量嵌入文档，单条失败时自动重试，仍然失败则抛出异常"""
        embeddings = []
        for i, text in enumerate(texts):
            try:
                label = f"第 {i+1}/{len(texts)} 条文本"
                result = self._embed_with_retry(text, max_retries=max_retries, label=label)
                embeddings.append(result)
            except Exception as e:
                raise Exception(f"第 {i+1}/{len(texts)} 条文本 Embedding 失败，已重试 {max_retries} 次：{e}") from e
        return embeddings

    def embed_query(self, text: str, max_retries: int = 2) -> List[float]:
        """嵌入查询文本，失败时自动重试"""
        return self._embed_with_retry(text, max_retries=max_retries, label="查询")


class KnowledgeBase:
    def __init__(
        self,
        name: str = "default",
        documents_path: Optional[str] = None,
        db_path: Optional[str] = None,
        embed_model: str = "text-embedding-qwen3-embedding-4b",
        llm_model: str = "qwen/qwen3.5-9b",
        chunk_size: int = 800,
        chunk_overlap: int = 100,
        lmstudio_base_url: str = "http://127.0.0.1:1234",
        api_key: str = None
    ):
        self.name = name

        # 动态路径推导：core/ 的父目录为项目根目录
        if documents_path is None:
            CORE_DIR = os.path.dirname(os.path.abspath(__file__))
            INSTALL_DIR = os.path.dirname(CORE_DIR)
            documents_path = os.path.join(INSTALL_DIR, "documents")

        self.documents_path = documents_path
        if db_path is None:
            db_path = os.path.join(documents_path, ".rag", "chroma_db")
        self.db_path = db_path
        self.embed_model = embed_model
        self.llm_model = llm_model
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.lmstudio_base_url = lmstudio_base_url
        self.api_key = api_key
        self.meta_path = os.path.join(documents_path, ".rag", "metadata.json")
        
        # 初始化 Embeddings（必须在 collection 之前）
        self.embeddings = LMStudioEmbeddings(model=embed_model, base_url=lmstudio_base_url, api_key=api_key)
        
        # 初始化 ChromaDB（使用自定义 Embedding 函数）
        os.makedirs(self.db_path, exist_ok=True)
        self.client = chromadb.PersistentClient(path=self.db_path)
        self.collection = self.client.get_or_create_collection(
            name=f"kb_{name}",
            embedding_function=self.embeddings
        )
        
        # 文本分块器
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        
        # 初始化 OCR 引擎（用于 Word 图片识别）
        self.ocr = None
        
        print(f"[OK] 知识库 '{name}' 初始化完成")
        print(f"   - LLM: {llm_model}")
        print(f"   - Embedding: {embed_model}")
        print(f"   - LM Studio: {lmstudio_base_url}")
    
    def _get_file_hash(self, filepath: str) -> str:
        hasher = hashlib.md5()
        with open(filepath, 'rb') as f:
            hasher.update(f.read())
        return hasher.hexdigest()
    
    def _load_metadata(self) -> Dict:
        if os.path.exists(self.meta_path):
            with open(self.meta_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {"documents": {}, "last_updated": None}
    
    def _save_metadata(self, metadata: Dict):
        os.makedirs(os.path.dirname(self.meta_path), exist_ok=True)
        with open(self.meta_path, 'w', encoding='utf-8') as f:
            json.dump(metadata, f, ensure_ascii=False, indent=2)
    
    def _get_ocr_txt_path(self, filepath: str) -> str:
        """根据原文件路径生成对应的 OCR 文本文件路径（不区分大小写扩展名）"""
        base, _ = os.path.splitext(filepath)
        return base + '_ocr.txt'

    def _load_pdf(self, filepath: str) -> str:
        """读取 PDF 文件（优先读取 OCR 文本）"""
        # 优先检查 OCR 文本文件
        ocr_txt_path = self._get_ocr_txt_path(filepath)
        if os.path.exists(ocr_txt_path):
            print(f"  [OCR] 使用 OCR 文本：{os.path.basename(ocr_txt_path)}")
            with open(ocr_txt_path, 'r', encoding='utf-8') as f:
                return f.read()

        # 否则直接读取 PDF 文本
        text = ""
        with pymupdf.open(filepath) as doc:
            for page in doc:
                text += page.get_text()
        return text
    
    def _load_docx(self, filepath: str) -> str:
        """读取 Word 文件（包含段落、表格和图片 OCR）"""
        try:
            from docx import Document
            from PIL import Image
            import io
            import numpy as np
            
            doc = Document(filepath)
            text = ""
            
            # 读取所有段落
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
            
            # 读取所有表格
            for i, table in enumerate(doc.tables, 1):
                text += f"\n=== 表格 {i} ===\n"
                for j, row in enumerate(table.rows):
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text += f"行{j}: {' | '.join(row_text)}\n"
                text += "\n"
            
            # 读取并 OCR 图片
            image_count = 0
            ocr_text = ""
            
            # 遍历文档中的所有关系（包括图片）
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        # 获取图片数据
                        image_blob = rel.target_part.blob
                        image = Image.open(io.BytesIO(image_blob))
                        
                        # 转换为 numpy 数组供 OCR 使用
                        image_array = np.array(image)
                        
                        # 初始化 OCR 引擎（如果还未初始化）
                        if self.ocr is None:
                            print(f"  [OCR] 初始化 OCR 引擎（首次使用）...")
                            import os
                            os.environ['PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK'] = 'True'
                            from paddleocr import PaddleOCR
                            self.ocr = PaddleOCR(use_textline_orientation=True, lang='ch')
                            print(f"  [OCR] OCR 引擎初始化完成")
                        
                        # OCR 识别
                        print(f"  [OCR] 识别图片 {image_count + 1}...")
                        ocr_result = self.ocr.ocr(image_array, cls=True)
                        
                        if ocr_result and ocr_result[0]:
                            ocr_text += f"\n[图片 {image_count + 1} 内容]\n"
                            for line in ocr_result[0]:
                                if line and len(line) > 1 and line[1][0]:
                                    ocr_text += line[1][0] + "\n"
                            ocr_text += "[/图片内容]\n"
                            print(f"  [OCR] 图片 {image_count + 1} 识别完成")
                        
                        image_count += 1
                    except Exception as e:
                        print(f"  [Warning] OCR 识别图片失败：{e}")
                        continue
            
            if image_count > 0:
                text += f"\n\n=== 图片 OCR 内容（共 {image_count} 张）===\n"
                text += ocr_text
            
            return text
        except Exception as e:
            print(f"[Warning] 读取 Word 文件失败：{e}")
            return ""
    
    def _load_excel(self, filepath: str) -> str:
        """读取 Excel 文件"""
        try:
            import pandas as pd
            # 尝试读取所有 sheet
            excel_file = pd.ExcelFile(filepath)
            text = ""
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(filepath, sheet_name=sheet_name)
                text += f"\n=== {sheet_name} ===\n"
                text += df.to_string() + "\n"
            return text
        except Exception as e:
            print(f"[Warning] 读取 Excel 文件失败：{e}")
            return ""
    
    def _load_csv(self, filepath: str) -> str:
        """读取 CSV 文件"""
        try:
            import pandas as pd
            df = pd.read_csv(filepath)
            return df.to_string()
        except Exception as e:
            print(f"[Warning] 读取 CSV 文件失败：{e}")
            return ""
    
    def _load_pptx(self, filepath: str) -> str:
        """读取 PowerPoint 文件"""
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text = ""
            for i, slide in enumerate(prs.slides, 1):
                text += f"\n=== 幻灯片 {i} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"[Warning] 读取 PPT 文件失败：{e}")
            return ""
    
    def _load_txt(self, filepath: str) -> str:
        """读取 TXT/Markdown 文件"""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            print(f"[Warning] 读取文本文件失败：{e}")
            return ""
    
    def _load_document(self, filepath: str) -> str:
        """根据文件类型加载文档"""
        lower_path = filepath.lower()
        
        if lower_path.endswith(".pdf"):
            return self._load_pdf(filepath)
        elif lower_path.endswith(".docx") or lower_path.endswith(".doc"):
            return self._load_docx(filepath)
        elif lower_path.endswith(".xlsx") or lower_path.endswith(".xls"):
            return self._load_excel(filepath)
        elif lower_path.endswith(".csv"):
            return self._load_csv(filepath)
        elif lower_path.endswith(".pptx") or lower_path.endswith(".ppt"):
            return self._load_pptx(filepath)
        elif lower_path.endswith((".txt", ".md", ".markdown")):
            return self._load_txt(filepath)
        else:
            # 其他格式尝试作为文本读取
            return self._load_txt(filepath)
    
    def _split_text(self, text: str) -> List[str]:
        """简单文本分块"""
        chunks = []
        start = 0
        while start < len(text):
            end = start + self.chunk_size
            chunk = text[start:end]
            # 尝试在句子边界切断
            if end < len(text):
                for sep in ["\n\n", "\n", "。", "！", "？"]:
                    pos = chunk.rfind(sep)
                    if pos > self.chunk_size // 2:
                        chunk = chunk[:pos + len(sep)]
                        break
            chunks.append(chunk.strip())
            start = end - self.chunk_overlap if end < len(text) else len(text)
        return chunks
    
    def _scan_documents(self) -> List[str]:
        """扫描文档目录 - 支持所有常见格式（包括 OCR 文本）"""
        documents = []
        
        # 定义支持的格式
        supported_extensions = [
            '.pdf', '.PDF',
            '.docx', '.DOCX', '.doc', '.DOC',
            '.txt', '.TXT', '.md', '.MD', '.markdown', '.MARKDOWN',
            '.xlsx', '.XLSX', '.xls', '.XLS', '.csv', '.CSV',
            '.pptx', '.PPTX', '.ppt', '.PPT'
        ]
        
        # 递归扫描所有子目录
        for root, dirs, files in os.walk(self.documents_path):
            # 跳过隐藏目录和.rag 目录
            dirs[:] = [d for d in dirs if not d.startswith('.')]
            
            for file in files:
                filepath = os.path.join(root, file)
                filename = file.lower()
                ext = os.path.splitext(file)[1]
                
                # 优先使用 OCR 文本文件（如果有）
                if ext.lower() == '.pdf':
                    # 检查是否有对应的 OCR 文本文件
                    ocr_txt = self._get_ocr_txt_path(filepath)
                    if os.path.exists(ocr_txt):
                        documents.append(ocr_txt)  # 使用 OCR 文本
                        continue
                
                if ext in supported_extensions:
                    documents.append(filepath)
        
        # 去重并返回
        unique_docs = list(set(documents))
        
        if unique_docs:
            print(f"\n[INFO] 扫描到 {len(unique_docs)} 个文档:")
            for doc in unique_docs:
                print(f"  - {os.path.basename(doc)}")
        
        return unique_docs
    
    def _detect_need_ocr(self, filepath: str) -> bool:
        """检测文件是否需要 OCR 处理"""
        ext = os.path.splitext(filepath)[1].lower()
        
        # PDF 文件检测
        if ext == '.pdf':
            try:
                doc = pymupdf.open(filepath)
                total_text = ""
                # 检查前 3 页（如果页数不足则检查全部），避免封面误判
                pages_to_check = min(3, len(doc))
                for page_idx in range(pages_to_check):
                    total_text += doc[page_idx].get_text()
                doc.close()

                # 如果检查页的文字总量很少，判定为图片 PDF
                if not total_text.strip() or len(total_text.strip()) < 100:
                    return True
            except Exception as e:
                print(f"  [WARN] 检测 PDF 失败：{e}")
                return False
        
        # TODO: Word 文档图片检测（需要额外逻辑）
        # TODO: PPT 文档图片检测（需要额外逻辑）
        
        return False
    
    def _auto_ocr(self, filepath: str) -> str:
        """自动对文件进行 OCR 处理
        
        Returns:
            OCR 文本文件路径，如果不需要 OCR 则返回原路径
        """
        ext = os.path.splitext(filepath)[1].lower()
        
        # 检查是否已有 OCR 文件
        ocr_txt = self._get_ocr_txt_path(filepath)
        if os.path.exists(ocr_txt):
            print(f"  [OCR] 已有 OCR 文件：{os.path.basename(ocr_txt)}")
            return ocr_txt

        # 检测是否需要 OCR
        if self._detect_need_ocr(filepath):
            print(f"  [OCR] 检测到图片 PDF，开始自动识别：{os.path.basename(filepath)}")
            
            try:
                # 导入 OCR 模块
                import sys
                sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
                from ocr_simple import ocr_pdf
                
                # 执行 OCR
                ocr_output = ocr_pdf(filepath)
                
                if ocr_output:
                    print(f"  [OCR] ✅ 识别完成：{os.path.basename(ocr_output)}")
                    return ocr_output
                else:
                    print(f"  [OCR] ⚠️ 识别失败，使用原文档")
                    return filepath
                    
            except Exception as e:
                print(f"  [OCR] ⚠️ OCR 处理失败：{e}，使用原文档")
                return filepath
        
        return filepath
    
    def add_documents(self, force: bool = False, auto_ocr: bool = True):
        """增量添加文档（自动检测并处理 OCR）
        
        Args:
            force: 是否强制重建索引
            auto_ocr: 是否自动检测并处理需要 OCR 的文件
        """
        print("\n[INFO] 检查文档更新...")
        
        metadata = self._load_metadata()
        indexed = metadata.get("documents", {})
        files = self._scan_documents()
        
        to_index = []
        ocr_processed = []
        
        for filepath in files:
            # 自动 OCR 处理
            if auto_ocr:
                processed_path = self._auto_ocr(filepath)
                if processed_path != filepath:
                    ocr_processed.append((filepath, processed_path))
                    filepath = processed_path
            
            file_hash = self._get_file_hash(filepath)
            if filepath not in indexed or indexed[filepath] != file_hash or force:
                to_index.append(filepath)
                print(f"  [NEW] {os.path.basename(filepath)}")
        
        if not to_index:
            print("  [OK] 文档已是最新")
            return
        
        # 加载并分块文档
        all_chunks = []
        all_ids = []
        all_metadatas = []
        
        for filepath in to_index:
            # 根据文件类型加载
            text = self._load_document(filepath)
            
            if not text:
                print(f"  [SKIP] {os.path.basename(filepath)} (无法读取)")
                continue
            
            chunks = self._split_text(text)
            for i, chunk in enumerate(chunks):
                chunk_id = f"{filepath}_{i}"
                all_chunks.append(chunk)
                all_ids.append(chunk_id)
                all_metadatas.append({
                    "source": os.path.relpath(filepath, self.documents_path),
                    "filepath": filepath,
                    "hash": self._get_file_hash(filepath)
                })
        
        # 批量嵌入并添加
        if all_chunks:
            print(f"  正在嵌入 {len(all_chunks)} 个文本块...")
            embeddings = self.embeddings.embed_documents(all_chunks)
            
            self.collection.add(
                documents=all_chunks,
                embeddings=embeddings,
                ids=all_ids,
                metadatas=all_metadatas
            )
            
            # 更新 metadata
            for filepath in to_index:
                indexed[filepath] = self._get_file_hash(filepath)
            metadata["documents"] = indexed
            metadata["last_updated"] = datetime.now().isoformat()
            self._save_metadata(metadata)
            
            # 输出 OCR 统计
            if ocr_processed:
                print(f"\n  [OCR] 本次自动处理 {len(ocr_processed)} 个文件:")
                for orig, ocr_file in ocr_processed:
                    print(f"    • {os.path.basename(orig)} → {os.path.basename(ocr_file)}")
            
            print(f"  [OK] 已索引 {len(to_index)} 个文件，{len(all_chunks)} 个文本块")
    
    def get_status(self) -> IndexStatus:
        """获取索引状态"""
        metadata = self._load_metadata()
        return IndexStatus(
            total_documents=len(metadata.get("documents", {})),
            total_chunks=self.collection.count(),
            last_updated=metadata.get("last_updated", "从未"),
            documents=list(metadata.get("documents", {}).keys())
        )
    
    def search(self, query: str, top_k: int = 5):
        """搜索相关文档"""
        self.add_documents()
        query_embedding = self.embeddings.embed_query(query)
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k
        )
        
        # 转换为 Document 格式
        docs = []
        if results['documents'] and results['documents'][0]:
            for i, content in enumerate(results['documents'][0]):
                metadata = results['metadatas'][0][i] if results['metadatas'] else {}
                docs.append(type('Document', (), {'page_content': content, 'metadata': metadata})())
        return docs
    
    def query_with_sources(self, question: str, top_k: int = 5, auto_refresh: bool = False, max_context_length: int = 8000) -> QueryResult:
        """带引用来源的查询
        
        Args:
            question: 用户问题
            top_k: 返回的文档块数量
            auto_refresh: 是否自动检查并刷新索引（默认 False）
            max_context_length: 最大上下文长度（token 数，默认 8000，避免超过模型上下文窗口）
        """
        start_time = time.time()  # ⏱️ 记录开始时间
        
        # 可选：自动刷新索引（仅在添加新文档后需要）
        if auto_refresh:
            self.add_documents()
        
        # 检索
        query_embedding = self.embeddings.embed_query(question)
        results = self.collection.query(
            query_embeddings=[query_embedding],
            n_results=top_k
        )
        
        if not results['documents'] or not results['documents'][0]:
            query_time = time.time() - start_time
            return QueryResult(
                answer="未在知识库中找到相关信息。",
                sources=[],
                query=question,
                timestamp=datetime.now().isoformat(),
                query_time=query_time
            )
        
        # 构建上下文 - 直接使用所有检索结果
        context_parts = []
        sources = []
        
        for i, content in enumerate(results['documents'][0]):
            preview = content[:600]
            context_parts.append(f"[来源{i+1}] {preview}")
            metadata = results['metadatas'][0][i] if results['metadatas'] else {}
            sources.append({
                "index": i + 1,
                "source": metadata.get("source", "未知"),
                "content_preview": preview[:150]
            })
        
        context = "\n\n".join(context_parts)
        
        # 调用 LM Studio 生成答案
        prompt_tokens = 0
        completion_tokens = 0
        total_tokens = 0
        
        try:
            llm_prompt = f"""基于以下文档内容回答问题：

{context}

问题：{question}

要求：
- 只回答一次，不要重复
- 简洁明了，列出关键要点
- 如果文档中没有相关信息，直接说"未在知识库中找到相关信息"

答案："""

            payload = {
                "model": self.llm_model,
                "prompt": llm_prompt,
                "temperature": 0,
                "max_tokens": 800,
                "stream": False
            }
            # 构建 headers，支持 API Key 认证
            headers = {}
            if self.api_key:
                headers["Authorization"] = f"Bearer {self.api_key}"
            
            response = requests.post(
                f"{self.lmstudio_base_url}/v1/completions",
                json=payload,
                headers=headers,
                timeout=60
            )
            if response.status_code == 200:
                result = response.json()
                answer = result.get("choices", [{}])[0].get("text", "").strip()
                
                # ✅ 提取 token 信息
                usage = result.get("usage", {})
                prompt_tokens = usage.get("prompt_tokens", 0)
                completion_tokens = usage.get("completion_tokens", 0)
                total_tokens = usage.get("total_tokens", 0)
                
                # 清理重复内容（qwen thinking 模式问题）
                if answer.count('\n\n答案：') > 1:
                    # 只保留最后一次答案
                    parts = answer.split('\n\n答案：')
                    answer = parts[-1].strip()
                if not answer:
                    answer = "根据检索到的文档，相关内容如下：\n\n" + context[:600]
            else:
                answer = f"LM Studio 调用失败：{response.status_code}"
        except Exception as e:
            answer = f"LM Studio 调用错误：{e}"
        
        # ⏱️ 计算耗时
        query_time = time.time() - start_time
        
        return QueryResult(
            answer=answer,
            sources=sources,
            query=question,
            timestamp=datetime.now().isoformat(),
            prompt_tokens=prompt_tokens,
            completion_tokens=completion_tokens,
            total_tokens=total_tokens,
            query_time=query_time
        )


if __name__ == "__main__":
    print("=" * 60)
    print("RAG 知识库工具测试 - LM Studio 简化版")
    print("=" * 60)
    
    kb = KnowledgeBase()
    status = kb.get_status()
    print(f"\n索引状态:")
    print(f"  文档数：{status.total_documents}")
    print(f"  文本块数：{status.total_chunks}")
    print(f"  最后更新：{status.last_updated}")
    
    print("\n" + "=" * 60)
    print("测试查询：2026 年 2 月电费单的计费电量是多少？")
    print("=" * 60)
    
    result = kb.query_with_sources("2026 年 2 月电费单的计费电量是多少？")
    print(f"\n{result.answer}")
    
    if result.sources:
        print("\n来源文档:")
        seen = set()
        for s in result.sources:
            if s['source'] not in seen:
                seen.add(s['source'])
                print(f"  - {s['source']}")
